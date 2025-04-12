"""File handlers for different file formats."""

import logging
from pathlib import Path
from typing import List, Optional, Union

import fitz  # PyMuPDF
import markdown
import openpyxl
from pptx import Presentation  # python-pptx
from docx import Document

# Configure logger
logger = logging.getLogger(__name__)


def read_text_file(file_path: Union[str, Path]) -> Optional[str]:
    """Read content from a text file."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read().strip()
    except UnicodeDecodeError:
        # Try with a different encoding if UTF-8 fails
        try:
            with open(file_path, "r", encoding="latin-1") as f:
                return f.read().strip()
        except Exception as e:
            logger.error(f"Error reading text file '{file_path}': {str(e)}")
            return None
    except Exception as e:
        logger.error(f"Error reading text file '{file_path}': {str(e)}")
        return None


def read_pdf_file(file_path: Union[str, Path]) -> Optional[str]:
    """Read content from a PDF file using PyMuPDF."""
    try:
        doc = fitz.open(file_path)
        text = ""
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text += page.get_text()
        doc.close()
        return text.strip()
    except Exception as e:
        logger.error(f"Error reading PDF file '{file_path}': {str(e)}")
        return None


def read_docx_file(file_path: Union[str, Path]) -> Optional[str]:
    """Read content from a DOCX file."""
    try:
        doc = Document(file_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text.strip()
    except Exception as e:
        logger.error(f"Error reading DOCX file '{file_path}': {str(e)}")
        return None


def read_excel_file(file_path: Union[str, Path]) -> Optional[str]:
    """Read content from an Excel file."""
    try:
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        text_parts = []

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_parts.append(f"Sheet: {sheet_name}")

            for row in sheet.iter_rows(values_only=True):
                # Filter out None values and convert to strings
                row_text = [str(cell) for cell in row if cell is not None]
                if row_text:  # Only add non-empty rows
                    text_parts.append(" | ".join(row_text))

        return "\n".join(text_parts).strip()
    except Exception as e:
        logger.error(f"Error reading Excel file '{file_path}': {str(e)}")
        return None


def read_pptx_file(file_path: Union[str, Path]) -> Optional[str]:
    """Read content from a PowerPoint file."""
    try:
        prs = Presentation(file_path)
        text_parts = []

        for slide_number, slide in enumerate(prs.slides, 1):
            text_parts.append(f"Slide {slide_number}:")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())

        return "\n\n".join(text_parts).strip()
    except Exception as e:
        logger.error(f"Error reading PowerPoint file '{file_path}': {str(e)}")
        return None


def read_markdown_file(file_path: Union[str, Path]) -> Optional[str]:
    """Read content from a Markdown file and convert to plain text."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            md_content = f.read()

        # Convert markdown to HTML
        html = markdown.markdown(md_content)

        # Simple HTML to text conversion (basic implementation)
        # In a real implementation, you might want to use a proper HTML parser
        text = (
            html.replace("<br>", "\n")
            .replace("</p>", "\n\n")
            .replace("</h1>", "\n\n")
            .replace("</h2>", "\n\n")
        )
        text = (
            text.replace("</h3>", "\n\n")
            .replace("</h4>", "\n\n")
            .replace("</h5>", "\n\n")
            .replace("</h6>", "\n\n")
        )
        text = text.replace("</li>", "\n").replace("</ul>", "\n").replace("</ol>", "\n")

        # Remove HTML tags
        import re

        text = re.sub(r"<[^>]+>", "", text)

        return text.strip()
    except Exception as e:
        logger.error(f"Error reading Markdown file '{file_path}': {str(e)}")
        return None


def read_file_content(file_path: Union[str, Path]) -> Optional[str]:
    """Read content from a file based on its extension."""
    file_path = Path(file_path)

    if not file_path.exists():
        logger.error(f"File does not exist: {file_path}")
        return None

    # Get file extension (lowercase)
    ext = file_path.suffix.lower()

    # Handle different file types
    if ext == ".pdf":
        return read_pdf_file(file_path)
    elif ext == ".docx":
        return read_docx_file(file_path)
    elif ext == ".xlsx" or ext == ".xls":
        return read_excel_file(file_path)
    elif ext == ".pptx" or ext == ".ppt":
        return read_pptx_file(file_path)
    elif ext == ".md" or ext == ".markdown":
        return read_markdown_file(file_path)
    else:
        # Default to text file
        return read_text_file(file_path)


def get_supported_extensions() -> List[str]:
    """Get a list of supported file extensions."""
    return [
        ".txt",
        ".csv",
        ".json",
        ".xml",
        ".html",
        ".htm",
        ".py",
        ".js",
        ".java",
        ".c",
        ".cpp",
        ".h",
        ".hpp",
        ".cs",
        ".rb",
        ".php",
        ".go",
        ".rs",
        ".swift",
        ".kt",
        ".scala",
        ".sh",
        ".bash",
        ".zsh",
        ".fish",
        ".md",
        ".markdown",
        ".rst",
        ".ini",
        ".cfg",
        ".conf",
        ".yaml",
        ".yml",
        ".toml",
        ".env",
        ".gitignore",
        ".pdf",
        ".docx",
        ".xlsx",
        ".xls",
        ".pptx",
        ".ppt",
    ]
